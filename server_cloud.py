from flask import Flask, request, Response
from flask_cors import CORS
import urllib.request, urllib.error, json, os, io

app = Flask(__name__)
CORS(app)

def extract_text(filename, data):
    ext = os.path.splitext(filename.lower())[1]
    try:
        if ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f'--- Slide {i} ---')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text.strip())
            return '\n'.join(lines)
        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f'--- Sheet: {sheet.title} ---')
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        lines.append('\t'.join(cells))
            return '\n'.join(lines)
        elif ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            lines = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ''
                if text.strip():
                    lines.append(f'--- Page {i} ---\n{text.strip()}')
            return '\n'.join(lines)
        elif ext in ('.docx', '.doc'):
            import docx
            doc = docx.Document(io.BytesIO(data))
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in ('.txt', '.csv', '.md'):
            return data.decode('utf-8', errors='replace')
        else:
            return None
    except Exception as e:
        raise RuntimeError(f'Could not parse {filename}: {e}')

@app.route('/api/claude', methods=['POST', 'OPTIONS'])
def claude_proxy():
    if request.method == 'OPTIONS':
        return Response(status=200)
    api_key = request.headers.get('X-Api-Key', '')
    body = request.get_data()
    try:
        req = urllib.request.Request(
            'https://api.anthropic.com/v1/messages',
            data=body,
            headers={'Content-Type': 'application/json', 'x-api-key': api_key, 'anthropic-version': '2023-06-01'},
            method='POST',
        )
        with urllib.request.urlopen(req, timeout=60) as r:
            result = r.read()
        return Response(result, status=200, content_type='application/json')
    except urllib.error.HTTPError as e:
        return Response(e.read(), status=e.code, content_type='application/json')
    except Exception as e:
        return Response(json.dumps({'error': str(e)}), status=500, content_type='application/json')

@app.route('/api/extract', methods=['POST', 'OPTIONS'])
def extract():
    if request.method == 'OPTIONS':
        return Response(status=200)
    try:
        f = request.files.get('file')
        if not f:
            return Response(json.dumps({'error': 'No file uploaded'}), status=400, content_type='application/json')
        data = f.read()
        text = extract_text(f.filename, data)
        if text is None:
            return Response(json.dumps({'error': 'Unsupported file type'}), status=400, content_type='application/json')
        return Response(json.dumps({'text': text, 'filename': f.filename, 'chars': len(text)}), status=200, content_type='application/json')
    except Exception as e:
        return Response(json.dumps({'error': str(e)}), status=400, content_type='application/json')

@app.route('/')
def health():
    return 'DAMI server running'

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 4321))
    app.run(host='0.0.0.0', port=port)
